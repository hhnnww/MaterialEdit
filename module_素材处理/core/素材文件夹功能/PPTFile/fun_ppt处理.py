from pathlib import Path

import pythoncom
import yaml
from pptx import Presentation
from pywintypes import com_error
from win32com.client import DispatchEx

from module_素材处理.core.setting import IMG_PATH


class PPTFile:
    def __init__(self, ppt_file: Path):
        self.ppt_file = ppt_file
        self.ppt_dir = ppt_file.parent / ppt_file.stem
        yaml_path = IMG_PATH / 'ppt-keyword.yaml'
        with open(yaml_path.as_posix(), 'r', encoding='utf-8') as p_f:
            p_f_dict = yaml.load(p_f, Loader=yaml.FullLoader)
            self.key_list = p_f_dict['replace_key']

    def replace_ppt_keyword(self):
        try:
            prs = Presentation(self.ppt_file.as_posix())
        except:
            pass
        else:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key in self.key_list:
                                if key[0] in run.text:
                                    run.text = run.text.replace(key[0], key[1])
            prs.save(self.ppt_file.as_posix())

    def fun_导出PNG(self):
        pythoncom.CoInitialize()
        ppt_app = DispatchEx('PowerPoint.Application')
        ppt_app.DisplayAlerts = 0
        try:
            ppt = ppt_app.Presentations.Open(self.ppt_file.as_posix())
        except:
            print('文件打开错误。')
            return None
        else:
            try:
                ppt.SaveAs(self.ppt_dir, 17)
            except com_error:
                print(f'错误文件，无法导出：{self.ppt_file.as_posix()}')
                ppt.Close()
                return None

            if self.ppt_file.suffix.lower() == '.ppt':
                ppt.SaveAs(self.ppt_file.with_suffix('.pptx'))
                self.ppt_file.unlink()

            ppt.Close()
            pythoncom.CoUninitialize()

            return 'ok'


if __name__ == '__main__':
    PPTFile(Path(r'X:\H000-H999\H0257\H0257\小夕素材(1).pptx')).fun_导出PNG()
